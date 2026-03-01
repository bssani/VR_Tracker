"""
VR Tracker Collision System — Presentation Generator
Run: python3 generate_ppt.py
Output: VR_Ergonomics_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── Color Palette ────────────────────────────────────────────────────────────
C_DARK    = RGBColor(0x1A, 0x1A, 0x2E)   # dark navy (background)
C_ACCENT  = RGBColor(0x00, 0xB4, 0xD8)   # cyan accent
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xE0, 0xF7, 0xFF)   # light blue text
C_YELLOW  = RGBColor(0xFF, 0xD1, 0x66)   # warning yellow
C_GREEN   = RGBColor(0x57, 0xCC, 0x99)   # safe green
C_RED     = RGBColor(0xFF, 0x6B, 0x6B)   # collision red
C_GRAY    = RGBColor(0xAA, 0xBB, 0xCC)   # subtle text
C_BOX_BG  = RGBColor(0x16, 0x21, 0x3E)   # slightly lighter navy

W = Inches(13.33)   # widescreen width
H = Inches(7.5)     # widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

BLANK = prs.slide_layouts[6]   # completely blank layout


# ─── Helper: add background rect ─────────────────────────────────────────────
def bg(slide, color=C_DARK):
    shape = slide.shapes.add_shape(1, 0, 0, W, H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.zorder = 0


# ─── Helper: add rect ────────────────────────────────────────────────────────
def rect(slide, left, top, width, height, fill_color, line_color=None, radius=False):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1.2)
    else:
        shape.line.fill.background()
    return shape


# ─── Helper: add text box ─────────────────────────────────────────────────────
def txt(slide, text, left, top, width, height,
        font_size=18, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, wrap=True, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return txBox


# ─── Helper: slide header bar ────────────────────────────────────────────────
def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.1), C_ACCENT)
    txt(slide, title,
        Inches(0.4), Inches(0.12), Inches(10), Inches(0.65),
        font_size=28, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
    if subtitle:
        txt(slide, subtitle,
            Inches(0.4), Inches(0.72), Inches(10), Inches(0.35),
            font_size=14, color=C_DARK, align=PP_ALIGN.LEFT)


# ─── Helper: accent line under header ────────────────────────────────────────
def divider(slide, top=Inches(1.1)):
    rect(slide, 0, top, W, Inches(0.04), C_ACCENT)


# ─── Helper: bullet block ────────────────────────────────────────────────────
def bullets(slide, items, left, top, width, height,
            font_size=16, color=C_WHITE, bullet_char="▸  "):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(6)
        run = p.add_run()
        run.text = bullet_char + item
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.name = "Calibri"


# ─── Helper: info card ────────────────────────────────────────────────────────
def card(slide, left, top, width, height, title, body_lines,
         title_color=C_ACCENT, body_color=C_LIGHT,
         bg_color=C_BOX_BG, border_color=C_ACCENT):
    rect(slide, left, top, width, height, bg_color, border_color)
    txt(slide, title,
        left + Inches(0.15), top + Inches(0.1),
        width - Inches(0.3), Inches(0.4),
        font_size=14, bold=True, color=title_color)
    body = "\n".join(body_lines)
    txt(slide, body,
        left + Inches(0.15), top + Inches(0.5),
        width - Inches(0.3), height - Inches(0.6),
        font_size=13, color=body_color, wrap=True)


# ─── Helper: flow step box ───────────────────────────────────────────────────
def flow_box(slide, left, top, width, height, step_num, title, desc,
             fill=C_BOX_BG):
    rect(slide, left, top, width, height, fill, C_ACCENT)
    # Step badge
    badge = slide.shapes.add_shape(9, left + Inches(0.15), top + Inches(0.12),
                                   Inches(0.45), Inches(0.45))
    badge.fill.solid()
    badge.fill.fore_color.rgb = C_ACCENT
    badge.line.fill.background()
    txt(slide, str(step_num),
        left + Inches(0.16), top + Inches(0.1),
        Inches(0.43), Inches(0.45),
        font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    txt(slide, title,
        left + Inches(0.7), top + Inches(0.12),
        width - Inches(0.85), Inches(0.4),
        font_size=14, bold=True, color=C_ACCENT)
    txt(slide, desc,
        left + Inches(0.15), top + Inches(0.6),
        width - Inches(0.3), height - Inches(0.7),
        font_size=12, color=C_LIGHT, wrap=True)


# ─── Helper: arrow between boxes ─────────────────────────────────────────────
def arrow(slide, left, top, width=Inches(0.3)):
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "▶"
    run.font.size = Pt(22)
    run.font.color.rgb = C_ACCENT


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)

# Accent bar top
rect(slide, 0, 0, W, Inches(0.08), C_ACCENT)
# Accent bar bottom
rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)
# Left accent stripe
rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)

txt(slide, "VR-Based Vehicle Ergonomics Testing",
    Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4),
    font_size=40, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
txt(slide, "How We're Using Virtual Reality to Design\nSafer, More Comfortable Vehicles",
    Inches(1.0), Inches(3.3), Inches(9), Inches(1.2),
    font_size=22, color=C_LIGHT, align=PP_ALIGN.LEFT)

# Divider line
rect(slide, Inches(1.0), Inches(3.1), Inches(6), Inches(0.05), C_ACCENT)

txt(slide, "Leadership Briefing  ·  2026",
    Inches(1.0), Inches(4.7), Inches(8), Inches(0.5),
    font_size=15, color=C_GRAY, align=PP_ALIGN.LEFT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — TABLE OF CONTENTS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Table of Contents")

items = [
    ("01", "The Problem",          "Why ergonomic issues persist in traditional vehicle development"),
    ("02", "The Human Factor",     "How the human body creates design blind spots"),
    ("03", "Our Solution",         "VR-based body tracking and collision measurement"),
    ("04", "How It Works",         "Four-step workflow from setup to insight"),
    ("05", "What We Measure",      "5 body points, 3 warning levels, real-time feedback"),
    ("06", "Accuracy & Calibration","Per-subject body modeling and measurement specs"),
    ("07", "The Data",             "What every session produces for engineering teams"),
    ("08", "Business Value",       "Cost, speed, and decision quality"),
    ("09", "Current Status",       "Where we are and what comes next"),
]

col_w = Inches(5.8)
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    left = Inches(0.5) + col * (col_w + Inches(0.5))
    top  = Inches(1.3) + row * Inches(1.1)
    rect(slide, left, top, col_w, Inches(0.95), C_BOX_BG, C_ACCENT)
    txt(slide, num,
        left + Inches(0.12), top + Inches(0.08),
        Inches(0.5), Inches(0.4),
        font_size=18, bold=True, color=C_ACCENT)
    txt(slide, title,
        left + Inches(0.65), top + Inches(0.06),
        col_w - Inches(0.8), Inches(0.38),
        font_size=15, bold=True, color=C_WHITE)
    txt(slide, desc,
        left + Inches(0.65), top + Inches(0.46),
        col_w - Inches(0.8), Inches(0.38),
        font_size=11, color=C_GRAY)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "The Problem",
           "Every car has a blind spot in the design process")

# Big stat
rect(slide, Inches(0.5), Inches(1.3), Inches(3.8), Inches(2.5), C_BOX_BG, C_RED)
txt(slide, "10×",
    Inches(0.7), Inches(1.45), Inches(3.4), Inches(1.2),
    font_size=72, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
txt(slide, "More expensive to fix ergonomic\nissues after physical prototyping",
    Inches(0.7), Inches(2.65), Inches(3.4), Inches(0.9),
    font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# Problem bullets
bullets(slide, [
    "Ergonomic issues are discovered too late — after manufacturing commitments are made",
    "Drivers bump their knee on the AC unit, scrape their shin on the door frame, or can't\n   comfortably enter/exit the vehicle",
    "Physical prototypes are expensive and slow to iterate — each design change costs weeks",
    "Subjective feedback ('it felt tight') is difficult to act on — engineers need numbers",
    "A standard-size mannequin cannot capture the full range of real driver body shapes and motions",
], Inches(4.8), Inches(1.35), Inches(8.0), Inches(5.5),
   font_size=15, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — THE HUMAN FACTOR
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "The Human Factor",
           "It's not a design problem — it's a human body problem")

quote_box = Inches(0.5)
rect(slide, quote_box, Inches(1.3), W - Inches(1.0), Inches(1.5),
     C_BOX_BG, C_ACCENT)
txt(slide,
    '"The gap between \'fits on paper\' and \'fits in real life\' '
    'is where discomfort and injury risks hide."',
    Inches(0.9), Inches(1.4), W - Inches(1.8), Inches(1.35),
    font_size=18, color=C_YELLOW, align=PP_ALIGN.CENTER)

cards_data = [
    ("Human Bodies Move in Arcs",
     ["Vehicles are measured in millimeters on a screen",
      "But a driver's knee swings in a 3D arc on entry/exit",
      "Static clearance measurements miss dynamic risk"]),
    ("Everyone Is Different",
     ["Leg length, hip width, and entry angle vary widely",
      "A design safe for 5th-percentile female may harm 95th-percentile male",
      "One 'average' mannequin can't represent your full customer base"]),
    ("Motion Reveals Problems",
     ["Critical collisions happen in motion, not at rest",
      "Entry and exit sequences are the highest-risk moments",
      "These cannot be tested with a static digital model"]),
]
for i, (title, lines) in enumerate(cards_data):
    left = Inches(0.5) + i * Inches(4.22)
    card(slide, left, Inches(3.0), Inches(4.0), Inches(3.9),
         title, lines)
    # Override body font inside card manually — redraw body text
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), Inches(3.6),
        Inches(3.7), Inches(3.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(5)
        run = p.add_run()
        run.text = "▸  " + line
        run.font.size = Pt(13)
        run.font.color.rgb = C_LIGHT
        run.font.name = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Our Solution",
           "Test the human body — before building the car")

pillars = [
    ("VR  Virtual Reality",
     C_ACCENT,
     ["The driver wears a VR headset and walks through\na full-scale digital vehicle interior",
      "No physical prototype required",
      "Any vehicle design can be tested on the same day it is created"]),
    ("Body Tracking",
     C_GREEN,
     ["5 sensors on the body (waist, both knees, both feet)\ncapture exact position in real time — 90 times per second",
      "A digital body model mirrors every movement",
      "Calibrated individually to each test subject's proportions"]),
    ("Collision Intelligence",
     C_YELLOW,
     ["The system automatically measures how close\nevery body part gets to every vehicle component",
      "Three alert levels: Safe / Warning / Collision",
      "Every measurement is logged to a data file for engineering review"]),
]
for i, (title, color, lines) in enumerate(pillars):
    left = Inches(0.5) + i * Inches(4.22)
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.6),
         C_BOX_BG, color)
    # Top color band
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(0.55), color)
    txt(slide, title,
        left + Inches(0.15), Inches(1.35),
        Inches(3.7), Inches(0.48),
        font_size=16, bold=True, color=C_DARK, align=PP_ALIGN.LEFT)
    txBox = slide.shapes.add_textbox(
        left + Inches(0.15), Inches(2.0),
        Inches(3.7), Inches(4.7))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, line in enumerate(lines):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = "▸  " + line
        run.font.size = Pt(14)
        run.font.color.rgb = C_LIGHT
        run.font.name = "Calibri"

txt(slide,
    "Real human movement.  Digital vehicle.  Data you can act on.",
    Inches(0.5), Inches(7.05), W - Inches(1.0), Inches(0.38),
    font_size=15, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — HOW IT WORKS (FLOWCHART)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "How It Works", "Four steps from setup to insight")

steps = [
    (1, "Strap On Sensors",
     "HMD headset + 5 Vive Trackers attached to waist, both knees, and both feet\n\n~2 min setup"),
    (2, "Calibrate",
     "Subject stands in T-Pose for 10 seconds. System measures individual body proportions.\n\n~10 sec"),
    (3, "Enter the Vehicle in VR",
     "Subject performs entry/exit sequence inside the digital vehicle. System tracks every movement.\n\n~5 min"),
    (4, "Review the Data",
     "CSV report generated: minimum clearances, warning events, collision log — ready for engineering review.\n\n~1 min"),
]
box_w = Inches(2.85)
box_h = Inches(4.5)
for i, (num, title, desc) in enumerate(steps):
    left = Inches(0.45) + i * (box_w + Inches(0.38))
    flow_box(slide, left, Inches(1.35), box_w, box_h, num, title, desc)
    if i < 3:
        arrow(slide,
              left + box_w + Inches(0.04),
              Inches(1.35) + box_h / 2 - Inches(0.2))

# Timeline bar
rect(slide, Inches(0.45), Inches(6.15), W - Inches(0.9), Inches(0.06), C_ACCENT)
txt(slide, "Total time per session:  ~10 minutes",
    Inches(0.45), Inches(6.25), W - Inches(0.9), Inches(0.45),
    font_size=14, color=C_ACCENT, align=PP_ALIGN.CENTER, bold=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — WHAT WE MEASURE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "What We Measure",
           "5 body points · 3 warning levels · real-time feedback")

# Left: body points table
rect(slide, Inches(0.4), Inches(1.3), Inches(6.5), Inches(5.5), C_BOX_BG, C_ACCENT)
txt(slide, "Body Tracking Points",
    Inches(0.55), Inches(1.4), Inches(6.2), Inches(0.45),
    font_size=16, bold=True, color=C_ACCENT)

rows = [
    ("Point",          "Sensor",     "Monitors Against",          C_GRAY),
    ("Waist / Hip",    "Tracker 1",  "Seat rail, center console",  C_WHITE),
    ("Left Knee",      "Tracker 2",  "Steering column, dashboard", C_WHITE),
    ("Right Knee",     "Tracker 3",  "Center console, glove box",  C_WHITE),
    ("Left Foot",      "Tracker 4",  "Dead pedal, footwell floor", C_WHITE),
    ("Right Foot",     "Tracker 5",  "Pedals, footwell",           C_WHITE),
]
col_lefts  = [Inches(0.55), Inches(2.05), Inches(3.55)]
col_widths = [Inches(1.45), Inches(1.45), Inches(3.25)]
for r, (p, s, m, color) in enumerate(rows):
    top = Inches(1.9) + r * Inches(0.75)
    if r == 0:
        for c, (cl, cw) in enumerate(zip(col_lefts, col_widths)):
            rect(slide, cl, top, cw, Inches(0.65), C_ACCENT)
    txt(slide, p, col_lefts[0], top + Inches(0.1), col_widths[0], Inches(0.5),
        font_size=13, bold=(r == 0), color=C_DARK if r == 0 else color)
    txt(slide, s, col_lefts[1], top + Inches(0.1), col_widths[1], Inches(0.5),
        font_size=13, bold=(r == 0), color=C_DARK if r == 0 else color)
    txt(slide, m, col_lefts[2], top + Inches(0.1), col_widths[2], Inches(0.5),
        font_size=13, bold=(r == 0), color=C_DARK if r == 0 else color)

# Right: warning levels
level_data = [
    ("SAFE",      "> 10 cm clearance",   "No action needed",            C_GREEN),
    ("WARNING",   "3 – 10 cm clearance", "Screen darkens · alert tone", C_YELLOW),
    ("COLLISION", "< 3 cm clearance",    "Screen darkens fully · impact\nsound · visual FX at contact point", C_RED),
]
for i, (level, dist, response, color) in enumerate(level_data):
    top = Inches(1.3) + i * Inches(1.7)
    rect(slide, Inches(7.2), top, Inches(5.7), Inches(1.55), C_BOX_BG, color)
    rect(slide, Inches(7.2), top, Inches(1.1), Inches(1.55), color)
    txt(slide, level,
        Inches(7.25), top + Inches(0.1), Inches(1.0), Inches(0.45),
        font_size=12, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    txt(slide, dist,
        Inches(8.45), top + Inches(0.08), Inches(4.3), Inches(0.42),
        font_size=14, bold=True, color=color)
    txt(slide, response,
        Inches(8.45), top + Inches(0.55), Inches(4.3), Inches(0.85),
        font_size=12, color=C_LIGHT)

txt(slide, "Thresholds are fully configurable per test session.",
    Inches(7.2), Inches(6.5), Inches(5.7), Inches(0.4),
    font_size=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — IMMERSIVE FEEDBACK
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Immersive Feedback in VR",
           "The driver feels the problem — not just sees it")

feedback = [
    ("Visual", "Screen Vignette",
     "As a body part approaches a vehicle component, the edges of the VR display darken progressively.\n\n"
     "Full darkness at collision level — the subject's natural instinct is to adjust without verbal instruction.",
     C_ACCENT),
    ("Audio", "Alert Tones",
     "A soft alert tone plays at warning distance.\n\n"
     "A sharp impact sound plays at collision distance.\n\n"
     "Real-world auditory instinct takes over — no coaching needed.",
     C_YELLOW),
    ("Spatial", "Impact Visual FX",
     "A particle burst appears precisely where contact would occur, anchoring the collision in 3D space.\n\n"
     "Subjects look toward the contact point naturally, reinforcing body awareness.",
     C_GREEN),
]
for i, (tag, title, body, color) in enumerate(feedback):
    left = Inches(0.4) + i * Inches(4.22)
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.7), C_BOX_BG, color)
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(0.5), color)
    txt(slide, f"{tag}  —  {title}",
        left + Inches(0.15), Inches(1.33),
        Inches(3.7), Inches(0.45),
        font_size=14, bold=True, color=C_DARK)
    txt(slide, body,
        left + Inches(0.15), Inches(1.95),
        Inches(3.7), Inches(4.85),
        font_size=13, color=C_LIGHT, wrap=True)

txt(slide,
    "Subjects react naturally. The data reflects how a real person behaves — not how they respond to verbal instructions.",
    Inches(0.4), Inches(7.1), W - Inches(0.8), Inches(0.35),
    font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — ACCURACY & CALIBRATION
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Accuracy & Calibration",
           "The system adapts to every body — not one average")

# Calibration process
rect(slide, Inches(0.4), Inches(1.3), Inches(5.5), Inches(5.6), C_BOX_BG, C_ACCENT)
txt(slide, "Calibration Process",
    Inches(0.6), Inches(1.4), Inches(5.1), Inches(0.45),
    font_size=16, bold=True, color=C_ACCENT)
cal_steps = [
    "Subject stands in T-Pose (arms extended, legs together)",
    "10-second hold — system automatically captures measurement",
    "4 segment lengths measured: Hip→Knee and Knee→Foot for both legs",
    "Height estimated from headset position (±5 cm, manual override available)",
    "Personal body model created — used for all measurements in this session",
    "Re-calibration available at any time without restarting the session",
]
bullets(slide, cal_steps,
        Inches(0.55), Inches(1.9), Inches(5.15), Inches(4.8),
        font_size=13, color=C_LIGHT)

# Accuracy specs
rect(slide, Inches(6.2), Inches(1.3), Inches(6.7), Inches(5.6), C_BOX_BG, C_GREEN)
txt(slide, "Accuracy Specifications",
    Inches(6.4), Inches(1.4), Inches(6.3), Inches(0.45),
    font_size=16, bold=True, color=C_GREEN)
specs = [
    ("Tracker hardware precision",  "±1–2 mm  (HTC Vive Tracker 3.0)"),
    ("Tracking update rate",         "90 times per second"),
    ("Collision check rate",         "30 times per second"),
    ("Minimum valid measurement",    "Any segment < 10 cm triggers\nre-calibration automatically"),
    ("Height estimation",           "HMD position × 0.92 correction\n(or manual entry)"),
    ("Warning threshold",            "10 cm  (configurable 0–50 cm)"),
    ("Collision threshold",          "3 cm   (configurable 0–20 cm)"),
]
for i, (label, value) in enumerate(specs):
    top = Inches(1.9) + i * Inches(0.7)
    txt(slide, label,
        Inches(6.35), top, Inches(3.1), Inches(0.6),
        font_size=12, color=C_GRAY)
    txt(slide, value,
        Inches(9.5), top, Inches(3.2), Inches(0.6),
        font_size=12, bold=True, color=C_WHITE)

txt(slide,
    "We are not testing a 'standard person.'  We are testing this person, in this body, in this vehicle.",
    Inches(0.4), Inches(7.1), W - Inches(0.8), Inches(0.35),
    font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — THE DATA
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "The Data We Collect",
           "Every session produces actionable engineering data")

# Summary report card
rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.6), C_BOX_BG, C_ACCENT)
rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(0.55), C_ACCENT)
txt(slide, "Session Summary Report   (1 row per session)",
    Inches(0.55), Inches(1.34), Inches(5.8), Inches(0.48),
    font_size=14, bold=True, color=C_DARK)
summary_items = [
    "Subject ID  &  Height",
    "Individual body segment measurements  (hip→knee, knee→foot)",
    "Minimum clearance recorded  — per vehicle component",
    "Total warning events  and  collision events",
    "Which body part triggered each event  +  timestamp",
    "Session date, duration, and test configuration",
]
bullets(slide, summary_items,
        Inches(0.55), Inches(1.95), Inches(5.7), Inches(4.7),
        font_size=13, color=C_LIGHT)

# Frame data card
rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.6), C_BOX_BG, C_YELLOW)
rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.55), C_YELLOW)
txt(slide, "Full Frame Data Log   (optional, 10 Hz)",
    Inches(7.05), Inches(1.34), Inches(5.7), Inches(0.48),
    font_size=14, bold=True, color=C_DARK)
frame_items = [
    "Position of all 5 body points  —  recorded every 100 ms",
    "Distance to every reference point  —  every frame",
    "Complete 3D motion path for post-session reconstruction",
    "Full timeline of warning state changes",
    "Enables statistical analysis across multiple subjects",
]
bullets(slide, frame_items,
        Inches(7.05), Inches(1.95), Inches(5.7), Inches(4.7),
        font_size=13, color=C_LIGHT)

txt(slide,
    "Output format:  CSV  —  opens directly in Excel · compatible with statistical analysis tools · no special software needed",
    Inches(0.4), Inches(7.08), W - Inches(0.8), Inches(0.38),
    font_size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 11 — REAL-WORLD TEST SCENARIO
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "A Typical Test Session",
           "Driver entry/exit sequence — step by step")

scenario_steps = [
    ("Approach", C_GREEN,
     "Subject A (178 cm, male) approaches the driver door and reaches for the handle.\n\nAll clearances: SAFE\nNo warnings triggered."),
    ("Entry — Step In", C_YELLOW,
     "Left knee swings toward dashboard during step-in motion.\n\n⚠  WARNING\nLeft Knee / Dashboard\nClearance: 7.3 cm"),
    ("Seated", C_GREEN,
     "Subject settles into the seat position.\n\nAll clearances return to SAFE.\nNo active warnings."),
    ("Exit — Step Out", C_RED,
     "Right knee catches center console during exit motion.\n\n🔴  COLLISION\nRight Knee / Center Console\nClearance: 1.1 cm"),
]
for i, (phase, color, text) in enumerate(scenario_steps):
    left = Inches(0.4) + i * Inches(3.18)
    rect(slide, left, Inches(1.3), Inches(3.0), Inches(4.6), C_BOX_BG, color)
    rect(slide, left, Inches(1.3), Inches(3.0), Inches(0.5), color)
    txt(slide, phase,
        left + Inches(0.15), Inches(1.32),
        Inches(2.7), Inches(0.45),
        font_size=14, bold=True, color=C_DARK)
    txt(slide, text,
        left + Inches(0.15), Inches(1.9),
        Inches(2.7), Inches(3.9),
        font_size=13, color=C_LIGHT, wrap=True)
    if i < 3:
        arrow(slide,
              left + Inches(3.0) + Inches(0.0),
              Inches(1.3) + Inches(4.6) / 2 - Inches(0.2))

# Post-test result box
rect(slide, Inches(0.4), Inches(6.1), W - Inches(0.8), Inches(1.2), C_BOX_BG, C_RED)
txt(slide,
    "Post-Test Output:   Minimum clearance: 1.1 cm  (Right Knee / Center Console)   "
    "→   Engineering recommendation: Widen center console cutout by 15 mm",
    Inches(0.6), Inches(6.2), W - Inches(1.2), Inches(0.95),
    font_size=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 12 — COMPARISON
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Why This Beats Traditional Methods",
           "Faster insights · Earlier in the process · Lower cost")

headers = ["", "Traditional Physical Testing", "Our VR System"]
rows_data = [
    ("When in the process",  "Post-prototype (too late)",       "Pre-prototype (early design)"),
    ("Cost per test",        "High — prototype + lab time",     "Low — virtual, fully repeatable"),
    ("Time to set up",       "Days to weeks",                   "Minutes"),
    ("Data output",          "Subjective verbal feedback",      "Quantified distance & collision data"),
    ("Design iteration",     "New physical prototype required", "Change digital model, retest same day"),
    ("Subject variation",    "Difficult to control",            "Same protocol every time, any body"),
    ("Collision detection",  "Visual inspection only",          "Millimeter-accurate, automated"),
]

header_tops  = [Inches(1.3), Inches(1.3), Inches(1.3)]
col_lefts_c  = [Inches(0.4), Inches(3.8), Inches(8.5)]
col_widths_c = [Inches(3.3), Inches(4.6), Inches(4.4)]

# Header row
for c, (hl, cl, cw) in enumerate(zip(headers, col_lefts_c, col_widths_c)):
    if hl:
        rect(slide, cl, Inches(1.3), cw, Inches(0.65),
             C_RED if c == 1 else C_GREEN)
        txt(slide, hl, cl + Inches(0.1), Inches(1.35),
            cw - Inches(0.2), Inches(0.55),
            font_size=14, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)

# Data rows
for r, (label, bad, good) in enumerate(rows_data):
    top = Inches(2.05) + r * Inches(0.72)
    bg_c = C_BOX_BG if r % 2 == 0 else RGBColor(0x12, 0x1C, 0x38)
    rect(slide, col_lefts_c[0], top, col_widths_c[0], Inches(0.68), bg_c)
    rect(slide, col_lefts_c[1], top, col_widths_c[1], Inches(0.68), bg_c)
    rect(slide, col_lefts_c[2], top, col_widths_c[2], Inches(0.68), bg_c)
    txt(slide, label, col_lefts_c[0] + Inches(0.1), top + Inches(0.1),
        col_widths_c[0] - Inches(0.2), Inches(0.55), font_size=13, bold=True, color=C_GRAY)
    txt(slide, "✗  " + bad, col_lefts_c[1] + Inches(0.1), top + Inches(0.1),
        col_widths_c[1] - Inches(0.2), Inches(0.55), font_size=13, color=C_RED)
    txt(slide, "✓  " + good, col_lefts_c[2] + Inches(0.1), top + Inches(0.1),
        col_widths_c[2] - Inches(0.2), Inches(0.55), font_size=13, color=C_GREEN)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 13 — SYSTEM ARCHITECTURE (SIMPLIFIED)
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "How It's Built",
           "Six modular systems working together — each independently configurable")

layers = [
    ("Hardware Layer",       "HMD headset  +  5 Vive Trackers (waist, knees, feet)",            C_ACCENT),
    ("Body Model",           "Digital skeleton mirrors every movement in real time  (90 FPS)",   C_LIGHT),
    ("Vehicle Reference",    "Marked target zones placed on dashboard, console, footwell, etc.", C_LIGHT),
    ("Collision Detector",   "Measures distance from each body point to each vehicle zone  (30 Hz)", C_LIGHT),
    ("Feedback System",      "Visual darkening  ·  Audio alert  ·  3D impact effect",           C_YELLOW),
    ("Data Logger",          "Saves every measurement to CSV for post-session engineering review",  C_GREEN),
]
box_h = Inches(0.72)
arrow_h = Inches(0.3)
total = len(layers)
start_top = Inches(1.35)
box_w = Inches(10.0)
left = (W - box_w) / 2

for i, (name, desc, color) in enumerate(layers):
    top = start_top + i * (box_h + arrow_h)
    rect(slide, left, top, box_w, box_h, C_BOX_BG, color)
    rect(slide, left, top, Inches(2.4), box_h, color)
    txt(slide, name,
        left + Inches(0.1), top + Inches(0.12),
        Inches(2.2), Inches(0.5),
        font_size=13, bold=True, color=C_DARK, align=PP_ALIGN.CENTER)
    txt(slide, desc,
        left + Inches(2.55), top + Inches(0.14),
        box_w - Inches(2.7), Inches(0.5),
        font_size=13, color=C_LIGHT)
    if i < total - 1:
        txt(slide, "▼",
            left + box_w / 2 - Inches(0.2),
            top + box_h,
            Inches(0.4), Inches(0.3),
            font_size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)

txt(slide,
    "All six components are modular — each can be updated, swapped, or tuned independently.",
    Inches(0.4), Inches(7.1), W - Inches(0.8), Inches(0.35),
    font_size=13, color=C_GRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 14 — BUSINESS VALUE
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Business Value",
           "Cost reduction · Speed to market · Data-driven decisions")

value_cards = [
    ("Cost Reduction", C_GREEN,
     "Catching an ergonomic issue in VR — before a physical prototype exists — eliminates "
     "costly late-stage design changes.\n\n"
     "Industry estimates: 10–100× cheaper to fix in digital vs. physical prototype stage.\n\n"
     "One avoided prototype revision can fund the entire VR system."),
    ("Speed to Market", C_ACCENT,
     "Testing does not wait for a physical vehicle to be built.\n\n"
     "Design teams can run ergonomic validation in parallel with CAD/CAM work, "
     "compressing the development timeline.\n\n"
     "Multiple design variants can be tested in a single day."),
    ("Data-Driven Decisions", C_YELLOW,
     "Instead of:\n"
     "\"It felt a bit tight\"\n\n"
     "Engineering teams receive:\n"
     "\"Right Knee / Center Console:\n"
     " min. clearance 1.1 cm at t=3.2s\n"
     " during door entry sequence\"\n\n"
     "Every design decision is backed by measurement."),
]
for i, (title, color, body) in enumerate(value_cards):
    left = Inches(0.4) + i * Inches(4.22)
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(5.8), C_BOX_BG, color)
    rect(slide, left, Inches(1.3), Inches(4.0), Inches(0.55), color)
    txt(slide, title,
        left + Inches(0.15), Inches(1.33),
        Inches(3.7), Inches(0.48),
        font_size=16, bold=True, color=C_DARK)
    txt(slide, body,
        left + Inches(0.15), Inches(2.0),
        Inches(3.7), Inches(4.9),
        font_size=13, color=C_LIGHT, wrap=True)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 15 — CURRENT STATUS
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header_bar(slide, "Current Status & Next Steps",
           "Core system complete — first live test session ready")

# Complete
rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(5.7), C_BOX_BG, C_GREEN)
rect(slide, Inches(0.4), Inches(1.3), Inches(6.1), Inches(0.55), C_GREEN)
txt(slide, "Complete  ✓",
    Inches(0.55), Inches(1.33), Inches(5.7), Inches(0.48),
    font_size=16, bold=True, color=C_DARK)
done_items = [
    "Body tracking — 5 points, 90 FPS real-time",
    "Dynamic body model (limb segments + joint spheres)",
    "Per-subject calibration (T-Pose, 10 sec)",
    "Collision detection & distance measurement",
    "Three-level warning system (Safe / Warning / Collision)",
    "Immersive feedback (vignette, audio, impact FX)",
    "Session data logging (summary + full frame CSV)",
    "Session state management (Idle → Calibrate → Test → Review)",
    "Desktop simulation mode (no headset required for QA/review)",
]
bullets(slide, done_items,
        Inches(0.55), Inches(1.95), Inches(5.7), Inches(4.8),
        font_size=13, color=C_LIGHT, bullet_char="✓  ")

# In Progress / Next
rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(5.7), C_BOX_BG, C_YELLOW)
rect(slide, Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.55), C_YELLOW)
txt(slide, "In Progress / Next Steps",
    Inches(7.05), Inches(1.33), Inches(5.7), Inches(0.48),
    font_size=16, bold=True, color=C_DARK)
next_items = [
    ("In Progress",  "On-screen HUD (live distance display)"),
    ("In Progress",  "Subject info input form"),
    ("In Progress",  "Visual materials for body model"),
    ("Next",         "First complete test level with vehicle mesh"),
    ("Next",         "First closed-loop session with real vehicle\nmodel and human subject"),
    ("Next",         "Statistical analysis across subject pool"),
    ("Future",       "Multi-vehicle comparative study"),
]
for i, (tag, item) in enumerate(next_items):
    top = Inches(1.95) + i * Inches(0.72)
    tag_color = C_YELLOW if tag == "In Progress" else \
                C_ACCENT if tag == "Next" else C_GRAY
    txt(slide, f"[{tag}]",
        Inches(7.05), top, Inches(1.5), Inches(0.6),
        font_size=11, bold=True, color=tag_color)
    txt(slide, item,
        Inches(8.55), top, Inches(4.2), Inches(0.6),
        font_size=13, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 16 — CLOSING
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
rect(slide, 0, 0, W, Inches(0.08), C_ACCENT)
rect(slide, 0, H - Inches(0.08), W, Inches(0.08), C_ACCENT)
rect(slide, 0, 0, Inches(0.08), H, C_ACCENT)

txt(slide, "Designed for Humans.\nValidated by Data.",
    Inches(1.0), Inches(1.0), Inches(11.3), Inches(1.9),
    font_size=42, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
rect(slide, Inches(1.0), Inches(2.95), Inches(6.0), Inches(0.05), C_ACCENT)

takeaways = [
    "This solves a real human problem — ergonomic failures affect real drivers, not just metrics.",
    "The system is measurement-first — every interaction becomes a data point, not an opinion.",
    "It is designed to scale — any vehicle model, any body size, any team can run a session.",
]
for i, t in enumerate(takeaways):
    top = Inches(3.1) + i * Inches(0.95)
    rect(slide, Inches(1.0), top, Inches(0.06), Inches(0.7), C_ACCENT)
    txt(slide, t,
        Inches(1.3), top, Inches(9.5), Inches(0.8),
        font_size=17, color=C_LIGHT, wrap=True)

txt(slide, "Questions?",
    Inches(1.0), Inches(6.2), Inches(5), Inches(0.75),
    font_size=32, bold=True, color=C_ACCENT)

# ─── Save ─────────────────────────────────────────────────────────────────────
output_path = "VR_Ergonomics_Presentation.pptx"
prs.save(output_path)
print(f"Saved: {output_path}  ({len(prs.slides)} slides)")
